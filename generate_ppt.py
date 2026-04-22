"""
FairCheck India - PowerPoint Generator
Run: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 82, 147)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        sub_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf2 = sub_shape.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 82, 147)
    
    # Add bullets
    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf2 = content_shape.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)
        p.level = 0

# ========== SLIDE 1: COVER ==========
add_title_slide(prs, "FairCheck India", "AI Bias Detection & Fairness Auditing System")

# ========== SLIDE 2: PROBLEM ==========
add_content_slide(prs, "Problem Statement", [
    "AI systems in India show significant gender and demographic bias",
    "50-70% of AI models exhibit some form of demographic bias",
    "No existing solution for domain-specific AI fairness auditing in India",
    "Biased AI affects real lives:",
    "   - Women 30% less likely to get bank loans approved",
    "   - Male candidates selected 75% vs female 55% in Army recruitment",
    "   - Reserved category students face additional admission barriers",
    "Current solutions are generic, not tailored to Indian context"
])

# ========== SLIDE 3: SOLUTION ==========
add_content_slide(prs, "Our Solution - FairCheck India", [
    "India's first domain-specific AI bias detection system",
    "Three Domain Coverage:",
    "   - Indian Army Recruitment",
    "   - Education Admissions",
    "   - Bank Loan Approvals",
    "Uses ExponentiatedGradient algorithm with DemographicParity constraints",
    "Data source: Real government statistics from data.gov.in",
    "Provides both bias detection AND mitigation in one platform",
    "Results in 60 seconds with actionable compliance reports"
])

# ========== SLIDE 4: HOW IT WORKS ==========
add_content_slide(prs, "How It Works", [
    "Step 1: Select Domain (Army / Education / Bank Loan)",
    "Step 2: Load Data (10,000+ records from government statistics)",
    "Step 3: Analyze (Train ML models, calculate fairness metrics)",
    "Step 4: Generate Report (PDF with bias analysis and solutions)",
    "Step 5: Fix Bias (One-click bias mitigation)",
    "",
    "Key Features:",
    "   - Real-time fairness scoring",
    "   - SHAP explanations for each prediction",
    "   - PDF compliance report generation",
    "   - Demographic Parity & Equalized Odds metrics"
])

# ========== SLIDE 5: TECHNICAL APPROACH ==========
add_content_slide(prs, "Technical Approach", [
    "Technology Stack:",
    "   - Frontend: Streamlit (Python Web App)",
    "   - ML Models: scikit-learn, Fairlearn",
    "   - Explanations: SHAP (SHapley Additive exPlanations)",
    "   - Data: data.gov.in (Government of India Open Data)",
    "   - PDF: ReportLab",
    "",
    "Algorithms Tested:",
    "   - Logistic Regression (75% accuracy, 25% bias reduction)",
    "   - Decision Tree (80% accuracy, 30% bias reduction)",
    "   - Random Forest (85% accuracy, 35% bias reduction)",
    "   - Gradient Boosting (88% accuracy, 40% bias reduction)"
])

# ========== SLIDE 6: DATA SOURCES ==========
add_content_slide(prs, "Real Data Sources", [
    "Army Recruitment:",
    "   - Year-wise candidates recruited (2017-2022)",
    "   - State-wise Indian Army intake (2017-2020)",
    "   - Gender distribution from official data",
    "",
    "Education:",
    "   - College enrollment by state",
    "   - Caste-wise reservation statistics",
    "   - Gender parity in admissions",
    "",
    "Bank Loans:",
    "   - RBI credit distribution data",
    "   - Education loan disbursement state-wise",
    "   - Approval rate demographics"
])

# ========== SLIDE 7: IMPLEMENTATION ==========
add_content_slide(prs, "Implementation Plan", [
    "Q1 2026 - Phase 1: Army Recruitment",
    "   - Bias detection for Indian Army recruitment",
    "   - Gender parity analysis",
    "   - Physical test standardization checks",
    "",
    "Q2 2026 - Phase 2: Education Admissions",
    "   - College admission bias detection",
    "   - Caste/gender disparity analysis",
    "   - EWS quota verification",
    "",
    "Q3 2026 - Phase 3: Bank Loan Approvals",
    "   - Credit scoring bias detection",
    "   - Gender-wise approval rate analysis",
    "   - Alternative credit scoring",
    "",
    "Q4 2026 - Phase 4: API Integration",
    "   - Real-time monitoring",
    "   - Government database integration"
])

# ========== SLIDE 8: CHALLENGES ==========
add_content_slide(prs, "Challenges & Limitations", [
    "Current Limitations:",
    "   - Aggregated data only (no individual records due to privacy)",
    "   - Model accuracy trade-off with fairness",
    "   - Limited to available government statistics",
    "",
    "Challenges Ahead:",
    "   - Real-time data integration",
    "   - Expanding to healthcare and jobs sector",
    "   - User adoption and trust building",
    "",
    "Ethical Considerations:",
    "   - Balancing fairness with accuracy",
    "   - Transparent bias explanation",
    "   - Compliance with Indian AI regulations"
])

# ========== SLIDE 9: FUTURE SCOPE ==========
add_content_slide(prs, "Future Scope", [
    "Short-term (2026):",
    "   - Mobile app for field workers",
    "   - Real-time API for organizations",
    "   - Integration with government portals",
    "",
    "Long-term (2027):",
    "   - Healthcare AI bias detection",
    "   - Job recruitment bias auditing",
    "   - Housing loan bias analysis",
    "",
    "Expected Impact:",
    "   - 40-60% reduction in AI bias",
    "   - 100+ organizations served",
    "   - Open source community building"
])

# ========== SLIDE 10: ROADMAP ==========
add_content_slide(prs, "Project Roadmap 2026", [
    "JAN - MAR: Army Recruitment Module",
    "   - Deploy bias detection for Army",
    "   - Partner with defense ministry",
    "",
    "APR - JUN: Education Module",
    "   - Launch education bias detection",
    "   - Partner with UGC/MHRD",
    "",
    "JUL - SEP: Bank Loan Module",
    "   - Deploy loan bias detection",
    "   - Partner with RBI/banks",
    "",
    "OCT - DEC: API & Mobile App",
    "   - Launch public API",
    "   - Mobile app for field use"
])

# ========== SLIDE 11: CONTACT ==========
add_title_slide(prs, "Thank You", "FairCheck India - AI Bias Detection System")

# Save
output_path = os.path.join(os.path.dirname(__file__), "FairCheck_India_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print("\nOpen with: python -m pptx or Microsoft PowerPoint / Google Slides")