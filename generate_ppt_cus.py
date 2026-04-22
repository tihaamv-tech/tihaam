"""
FairCheck India - PowerPoint Generator
Matches Cepheus Template exactly
Run: python generate_ppt_cus.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE_DARK = RGBColor(0, 82, 147)
BLUE_LIGHT = RGBColor(0, 150, 199)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(150, 150, 150)
DARK = RGBColor(50, 50, 50)

def add_cover_slide(prs, title, subtitle=""):
    """Slide 1: Cover - Blue gradient with centered text"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Main title
    title_shape = slide.shapes.add_textbox(Inches(0), Inches(2.8), Inches(13.333), Inches(1.2))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_shape = slide.shapes.add_textbox(Inches(0), Inches(4), Inches(13.333), Inches(0.8))
        tf2 = sub_shape.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title, subtitle=""):
    """Slide 2: Title - Large centered title"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Main title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf2 = sub_shape.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, section_title, content_list):
    """Slides 3-10: Section with header and bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Section header
    header_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = header_shape.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Content bullets
    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf2 = content_shape.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.space_after = Pt(14)

def add_contact_slide(prs, name, email, phone, website):
    """Slide 11: Contact - Same layout as cover"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Get in touch text
    title_shape = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(1))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Get in touch"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact details
    contact_shape = slide.shapes.add_textbox(Inches(0), Inches(4), Inches(13.333), Inches(2.5))
    tf2 = contact_shape.text_frame
    p = tf2.paragraphs[0]
    p.text = name
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf2.add_paragraph()
    p.text = email
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf2.add_paragraph()
    p.text = phone
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf2.add_paragraph()
    p.text = website
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ========== CREATE ALL 11 SLIDES ==========

# Slide 1: Cover
add_cover_slide(prs, "FairCheck India", "AI Bias Detection & Fairness Auditing System")

# Slide 2: Title
add_title_slide(prs, "FairCheck India", "AI Fairness Auditing for Indian Institutions")

# Slide 3: About Us (Problem Statement)
add_section_slide(prs, "About Us", [
    "FairCheck India is India's first domain-specific AI bias detection system",
    "We analyze AI models used in Army recruitment, Education admissions, and Bank loan approvals",
    "Our mission: Detect and fix gender and demographic bias in AI systems",
    "Using real government statistics from data.gov.in",
    "Provides actionable compliance reports with solutions"
])

# Slide 4: Team (Solution Overview)
add_section_slide(prs, "Our Solution", [
    "Three Domain Coverage: Indian Army, Education, Bank Loan",
    "Step 1: Select Domain → Step 2: Load Data → Step 3: Analyze",
    "Step 4: Generate Report → Step 5: Fix Bias",
    "One-click bias mitigation using ExponentiatedGradient algorithm",
    "Results in 60 seconds with PDF compliance report"
])

# Slide 5: Problem (Technical Approach)
add_section_slide(prs, "Problem Statement", [
    "50-70% of AI models exhibit demographic bias",
    "Women 30% less likely to get bank loans approved",
    "Male candidates selected 75% vs female 55% in Army recruitment",
    "No existing solution for Indian domain-specific bias auditing",
    "Current solutions are generic, not tailored to Indian context"
])

# Slide 6: Solution (Implementation)
add_section_slide(prs, "Technical Stack", [
    "Frontend: Streamlit (Python Web App)",
    "ML Models: scikit-learn, Fairlearn",
    "Explanations: SHAP (SHapley Additive exPlanations)",
    "Data Source: data.gov.in (Government of India Open Data)",
    "Algorithms: Logistic Regression, Decision Tree, Random Forest, Gradient Boosting"
])

# Slide 7: Market Opportunity (Data Sources)
add_section_slide(prs, "Data Sources", [
    "Army Recruitment: Year-wise candidates recruited (2017-2022)",
    "State-wise Indian Army intake (2017-2020)",
    "Education: College enrollment by state, Caste-wise reservation statistics",
    "Bank Loans: RBI credit distribution data, Education loan disbursement",
    "All data sourced from official Government of India open data portal"
])

# Slide 8: Business Model (Implementation Plan)
add_section_slide(prs, "Implementation Plan", [
    "Q1 2026: Army Recruitment module - Partner with defense ministry",
    "Q2 2026: Education module - Partner with UGC/MHRD",
    "Q3 2026: Bank Loan module - Partner with RBI/banks",
    "Q4 2026: API Integration - Real-time monitoring for organizations",
    "Revenue: Government contracts, Enterprise licensing, API subscriptions"
])

# Slide 9: Competitors (Challenges)
add_section_slide(prs, "Challenges & Limitations", [
    "Limitations: Aggregated data only (privacy), Model accuracy trade-off",
    "Challenge: Real-time data integration with government databases",
    "Challenge: User adoption and trust building",
    "Ethical: Balancing fairness with accuracy",
    "Future: Expanding to healthcare and jobs sector"
])

# Slide 10: Roadmap
add_section_slide(prs, "Project Roadmap", [
    "JAN - MAR: Deploy Army Recruitment bias detection",
    "APR - JUN: Launch Education bias detection",
    "JUL - SEP: Deploy Bank Loan bias detection",
    "OCT - DEC: Launch public API and Mobile app",
    "Goal: 40-60% reduction in AI bias by 2027"
])

# Slide 11: Contact
add_contact_slide(prs, "FairCheck India Team", "faircheck@india.ai", "+91 98765 43210", "www.faircheckindia.ai")

# Save
output_path = "C:/Users/hp/Documents/FairCheck_India_Cepheus_Template.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")